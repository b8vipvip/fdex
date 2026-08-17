from __future__ import annotations

import base64
import io

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.document_service import extract_document_text, prepare_documents


def test_plain_text_document_is_injected_into_prompt() -> None:
    raw = "客户名称：星河科技\n预算：20 万\n目标：9 月上线".encode("utf-8")
    prepared = prepare_documents(
        "请分析附件",
        [{"name": "requirements.txt", "mime_type": "text/plain", "data": base64.b64encode(raw).decode()}],
    )
    assert prepared.extracted_count == 1
    assert "星河科技" in prepared.prompt
    assert "20 万" in prepared.prompt


def test_docx_extraction_reads_paragraphs() -> None:
    stream = io.BytesIO()
    doc = Document()
    doc.add_paragraph("FDEX 文档解析测试")
    doc.add_paragraph("关键风险：交付时间不足")
    doc.save(stream)
    text = extract_document_text(
        name="plan.docx",
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        raw=stream.getvalue(),
    )
    assert "FDEX 文档解析测试" in text
    assert "交付时间不足" in text


def test_xlsx_extraction_reads_cells() -> None:
    stream = io.BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "预算"
    sheet.append(["项目", "金额"])
    sheet.append(["服务器", 8888])
    workbook.save(stream)
    workbook.close()
    text = extract_document_text(
        name="budget.xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        raw=stream.getvalue(),
    )
    assert "预算" in text
    assert "服务器" in text
    assert "8888" in text


def test_pptx_extraction_reads_slide_text() -> None:
    stream = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "市场计划"
    slide.placeholders[1].text = "先验证付费转化，再扩大投放"
    presentation.save(stream)
    text = extract_document_text(
        name="market.pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        raw=stream.getvalue(),
    )
    assert "市场计划" in text
    assert "扩大投放" in text


def test_unsupported_document_is_explicitly_marked_unread() -> None:
    prepared = prepare_documents(
        "分析附件",
        [{"name": "archive.bin", "mime_type": "application/octet-stream", "data": base64.b64encode(b"abc").decode()}],
    )
    assert prepared.extracted_count == 0
    assert "正文提取失败" in prepared.prompt
