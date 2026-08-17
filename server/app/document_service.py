from __future__ import annotations

import base64
import binascii
import io
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

MAX_DOCUMENT_BYTES = 8 * 1024 * 1024
MAX_DOCUMENTS = 3
MAX_EXTRACTED_CHARS_PER_DOCUMENT = 18_000
MAX_EXTRACTED_CHARS_TOTAL = 36_000

_TEXT_SUFFIXES = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl", ".xml",
    ".log", ".conf", ".ini", ".yaml", ".yml", ".env", ".properties", ".sql",
    ".py", ".kt", ".kts", ".java", ".js", ".ts", ".tsx", ".jsx", ".html", ".css",
}


@dataclass
class DocumentPreparation:
    prompt: str
    extracted_count: int
    notes: list[str]


def prepare_documents(prompt: str, documents: list[dict[str, str]]) -> DocumentPreparation:
    sections: list[str] = []
    notes: list[str] = []
    used_chars = 0
    extracted_count = 0

    for item in documents[:MAX_DOCUMENTS]:
        name = str(item.get("name") or "附件").strip()[:240] or "附件"
        mime = str(item.get("mime_type") or "application/octet-stream").strip().lower()[:120]
        encoded = str(item.get("data") or "").strip()
        try:
            raw = base64.b64decode(encoded, validate=True)
        except (binascii.Error, ValueError):
            notes.append(f"文件《{name}》Base64 无效，未读取内容。")
            continue
        if not raw:
            notes.append(f"文件《{name}》为空，未读取内容。")
            continue
        if len(raw) > MAX_DOCUMENT_BYTES:
            notes.append(f"文件《{name}》超过 8 MB，未读取内容。")
            continue

        try:
            text = extract_document_text(name=name, mime_type=mime, raw=raw)
        except Exception as exc:
            notes.append(f"文件《{name}》正文提取失败：{exc.__class__.__name__}。")
            continue
        text = _clean(text)
        if not text:
            notes.append(f"文件《{name}》没有提取到可读文字，未假装已经读取。")
            continue

        remaining = MAX_EXTRACTED_CHARS_TOTAL - used_chars
        if remaining <= 0:
            notes.append("文档正文总量已达到本次分析上限，其余内容未继续送入模型。")
            break
        clipped = text[: min(MAX_EXTRACTED_CHARS_PER_DOCUMENT, remaining)]
        used_chars += len(clipped)
        extracted_count += 1
        suffix_note = "\n[正文已截断]" if len(clipped) < len(text) else ""
        sections.append(f"\n\n--- 文件正文：{name}（{mime}）---\n{clipped}{suffix_note}\n--- 文件正文结束 ---")

    if documents and len(documents) > MAX_DOCUMENTS:
        notes.append(f"一次最多解析 {MAX_DOCUMENTS} 份文档，后续文档未读取。")

    combined = (prompt or "").rstrip()
    if sections:
        combined += "\n\n以下内容是 FDEX 从用户实际附件中提取的正文，请只基于实际提取到的内容判断：" + "".join(sections)
    if notes:
        combined += "\n\n文件提取说明：\n" + "\n".join(f"- {note}" for note in notes)
    return DocumentPreparation(prompt=combined, extracted_count=extracted_count, notes=notes)


def extract_document_text(*, name: str, mime_type: str, raw: bytes) -> str:
    suffix = Path(name).suffix.lower()
    mime = (mime_type or "").lower()

    if suffix in _TEXT_SUFFIXES or mime.startswith("text/") or mime in {
        "application/json", "application/xml", "application/x-yaml", "application/yaml",
    }:
        return _decode_text(raw)
    if suffix == ".pdf" or mime == "application/pdf":
        reader = PdfReader(io.BytesIO(raw), strict=False)
        page_count = min(len(reader.pages), 80)
        return "\n".join((reader.pages[index].extract_text() or "") for index in range(page_count))
    if suffix == ".docx" or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(io.BytesIO(raw))
        parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
        for table in doc.tables[:30]:
            for row in table.rows[:200]:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    if suffix in {".xlsx", ".xlsm"} or "spreadsheetml" in mime:
        workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts: list[str] = []
        try:
            for sheet in workbook.worksheets[:12]:
                parts.append(f"[工作表 {sheet.title}]")
                for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if row_index >= 500:
                        break
                    values = [str(value) for value in row[:40] if value is not None]
                    if values:
                        parts.append("\t".join(values))
                    if sum(len(x) for x in parts) > MAX_EXTRACTED_CHARS_PER_DOCUMENT * 2:
                        break
        finally:
            workbook.close()
        return "\n".join(parts)
    if suffix == ".pptx" or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        presentation = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        slide_count = min(len(presentation.slides), 80)
        for index in range(slide_count):
            slide = presentation.slides[index]
            parts.append(f"[幻灯片 {index + 1}]")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    parts.append(str(text))
        return "\n".join(parts)

    raise ValueError("unsupported_document_type")


def _decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "big5", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _clean(text: str) -> str:
    value = (text or "").replace("\x00", "")
    lines = [line.rstrip() for line in value.splitlines()]
    return "\n".join(lines).strip()
